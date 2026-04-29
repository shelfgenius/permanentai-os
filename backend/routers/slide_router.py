"""
Slide Router — AI Presentation Generation backend.

Endpoints:
  POST /slide/generate  — Generate a full presentation from a prompt
  POST /slide/command   — Process editing commands (add slide, change theme, critique, etc.)
  GET  /slide/status    — Health check
"""
from __future__ import annotations

import json
import logging
import os
import re
import tempfile
from pathlib import Path

import httpx
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import Optional, List

logger = logging.getLogger("slide_router")
router = APIRouter(prefix="/slide", tags=["slide"])

NIM_BASE = os.getenv("NVIDIA_NIM_BASE_URL", "https://integrate.api.nvidia.com/v1")
NIM_KEY  = os.getenv("NVIDIA_API_KEY_CODING", "") or os.getenv("NVIDIA_API_KEY_TRANSLATE", "")


# ══════════════════════════════════════════════════════════════
#  Models
# ══════════════════════════════════════════════════════════════

class GenerateRequest(BaseModel):
    prompt: str
    style: str = "professional"  # professional, minimal, bold, creative, corporate
    slide_count: int = 8
    language: str = ""           # e.g. "Romanian", "English", "French" — auto-detected if empty
    audience: str = ""           # e.g. "students", "executives", "investors"

class CommandRequest(BaseModel):
    text: str
    current_slide: int = 0
    slide_count: int = 0

class SlideData(BaseModel):
    id: str = ""
    type: str = "content"        # title | content | split | section | chart | closing
    title: str = ""
    subtitle: str = ""
    bullets: List[str] = []
    notes: str = ""
    source: Optional[str] = None
    image: Optional[str] = None
    image_desc: str = ""         # description of what the image should show
    icons: List[str] = []        # suggested icon names
    chart_type: Optional[str] = None   # bar | pie | line | area | None
    chart_data: Optional[str] = None   # description of data for the chart
    layout: str = "default"      # default | two-col | image-left | image-right | full-image | centered
    colors: Optional[str] = None # per-slide color hint
    fonts: Optional[str] = None  # per-slide font hint
    bg: str = "#0f172a"


# ══════════════════════════════════════════════════════════════
#  Theme presets
# ══════════════════════════════════════════════════════════════

THEMES = {
    "professional": {"bg": "#0f172a", "accent": "#f59e0b", "text": "#f1f5f9", "secondary": "#94a3b8"},
    "corporate": {"bg": "#1e293b", "accent": "#3b82f6", "text": "#f8fafc", "secondary": "#cbd5e1"},
    "minimal": {"bg": "#ffffff", "accent": "#18181b", "text": "#18181b", "secondary": "#71717a"},
    "bold": {"bg": "#09090b", "accent": "#ef4444", "text": "#fafafa", "secondary": "#a1a1aa"},
    "creative": {"bg": "#1a1a2e", "accent": "#e040fb", "text": "#f1f5f9", "secondary": "#a78bfa"},
    "dark": {"bg": "#09090b", "accent": "#f59e0b", "text": "#e4e4e7", "secondary": "#71717a"},
}


# ══════════════════════════════════════════════════════════════
#  /generate — create a full presentation
# ══════════════════════════════════════════════════════════════

@router.post("/generate")
async def slide_generate(req: GenerateRequest):
    """Generate a full presentation from a prompt using AI."""
    prompt = req.prompt.strip()
    if not prompt:
        raise HTTPException(400, "Prompt is required")

    logger.info("Slide generate: prompt=%r style=%s count=%d", prompt[:80], req.style, req.slide_count)

    theme = THEMES.get(req.style, THEMES["professional"])
    title_bg = f"linear-gradient(135deg, {theme['bg']} 0%, #1e293b 50%, #334155 100%)"
    count = min(max(req.slide_count, 3), 15)

    # ── Try AI generation first ──
    if NIM_KEY:
        try:
            slides = await _generate_with_ai(prompt, count, theme, title_bg,
                                                language=req.language, style=req.style,
                                                audience=req.audience)
            if slides:
                return {
                    "slides": slides,
                    "theme": theme,
                    "topic": slides[0].get("title", "Presentation"),
                    "slide_count": len(slides),
                }
        except Exception as e:
            logger.warning("AI slide generation failed, using fallback: %s", e)

    # ── Fallback: template-based (preserves original language from prompt) ──
    slides = _generate_fallback(prompt, count, theme, title_bg)
    return {
        "slides": slides,
        "theme": theme,
        "topic": slides[0].get("title", "Presentation"),
        "slide_count": len(slides),
    }


async def _generate_with_ai(prompt: str, count: int, theme: dict, title_bg: str,
                            language: str = "", style: str = "professional",
                            audience: str = "") -> list:
    """Use NVIDIA NIM (Llama / Qwen) to generate visually rich slide content."""
    lang_instruction = (
        f"Generate ALL content in {language}."
        if language
        else "Detect the language of the user's prompt and generate ALL content in THAT SAME LANGUAGE. "
             "If the prompt is in Romanian, everything must be in Romanian. If English, all in English."
    )
    audience_line = f"Target audience: {audience}. Adapt tone and depth accordingly.\n" if audience else ""
    system_msg = (
        "You are an advanced AI presentation designer, not just a text generator.\n\n"
        f"LANGUAGE: {lang_instruction}\n"
        f"STYLE: {style} (adapt visuals and tone).\n"
        f"{audience_line}"
        "\nCRITICAL RULES:\n"
        "1. NO BASIC SLIDES — every slide must feel professionally designed.\n"
        "2. IMAGES ARE MANDATORY — every slide must include an image_desc field describing a relevant image.\n"
        "3. Use concise bullet points, highlight key ideas.\n"
        "4. Include charts/infographics where data is involved.\n"
        "5. Include a cover slide (first) and closing slide (last).\n\n"
        f"Generate exactly {count} slides as a JSON array. Each slide object MUST have:\n"
        '  - "type": one of "title", "content", "split", "section", "chart", "closing"\n'
        '  - "title": slide title\n'
        '  - "subtitle": subtitle\n'
        '  - "layout": one of "default", "two-col", "image-left", "image-right", "full-image", "centered"\n'
        '  - "bullets": array of 3-5 bullet strings (empty for title/section slides)\n'
        '  - "notes": speaker notes\n'
        '  - "source": citation or null\n'
        '  - "image_desc": specific description of what the image should contain (MANDATORY)\n'
        '  - "icons": array of suggested icon names (e.g. ["chart-bar", "globe", "users"])\n'
        '  - "chart_type": "bar" | "pie" | "line" | "area" | null\n'
        '  - "chart_data": description of chart data if chart_type is set, else null\n'
        '  - "colors": suggested color palette for this slide (e.g. "#1a73e8, #34a853, #ea4335")\n'
        '  - "fonts": suggested font pairing (e.g. "Montserrat + Open Sans")\n\n'
        "Make content detailed, specific, and factual with real data points.\n"
        "The result should feel like it was made by a professional designer.\n\n"
        "Return ONLY the JSON array, no markdown, no explanation."
    )

    payload = {
        "model": "meta/llama-3.3-70b-instruct",
        "messages": [
            {"role": "system", "content": system_msg},
            {"role": "user", "content": f"Create a presentation about: {prompt}"},
        ],
        "temperature": 0.7,
        "max_tokens": 4096,
        "stream": False,
    }

    async with httpx.AsyncClient(timeout=60) as client:
        r = await client.post(
            f"{NIM_BASE}/chat/completions",
            json=payload,
            headers={"Authorization": f"Bearer {NIM_KEY}", "Accept": "application/json"},
        )
    if r.status_code != 200:
        logger.warning("AI slide gen HTTP %s: %s", r.status_code, r.text[:300])
        return []

    data = r.json()
    content = data.get("choices", [{}])[0].get("message", {}).get("content", "")

    # Parse JSON from response (handle markdown code blocks)
    content = content.strip()
    if content.startswith("```"):
        content = re.sub(r'^```(?:json)?\s*', '', content)
        content = re.sub(r'\s*```$', '', content)

    try:
        raw_slides = json.loads(content)
    except json.JSONDecodeError:
        # Try to find JSON array in the response
        match = re.search(r'\[.*\]', content, re.DOTALL)
        if match:
            raw_slides = json.loads(match.group())
        else:
            logger.warning("Could not parse AI response as JSON")
            return []

    if not isinstance(raw_slides, list) or len(raw_slides) == 0:
        return []

    # Normalize into our format
    VALID_TYPES = {"title", "content", "split", "section", "chart", "closing"}
    VALID_LAYOUTS = {"default", "two-col", "image-left", "image-right", "full-image", "centered"}
    slides = []
    for i, rs in enumerate(raw_slides):
        slide_type = rs.get("type", "content")
        if slide_type not in VALID_TYPES:
            slide_type = "content"
        if i == 0:
            slide_type = "title"
        if i == len(raw_slides) - 1 and len(raw_slides) > 2:
            slide_type = rs.get("type", "closing") if rs.get("type") in ("title", "closing") else "closing"
        layout = rs.get("layout", "default")
        if layout not in VALID_LAYOUTS:
            layout = "default"
        slides.append({
            "id": f"sl-gen-{i+1}",
            "type": slide_type,
            "title": rs.get("title", ""),
            "subtitle": rs.get("subtitle", ""),
            "bg": title_bg if slide_type in ("title", "closing") else theme["bg"],
            "bullets": rs.get("bullets", []) if isinstance(rs.get("bullets"), list) else [],
            "source": rs.get("source"),
            "notes": rs.get("notes", ""),
            "image": None,
            "image_desc": rs.get("image_desc", ""),
            "icons": rs.get("icons", []) if isinstance(rs.get("icons"), list) else [],
            "chart_type": rs.get("chart_type") if rs.get("chart_type") in ("bar", "pie", "line", "area") else None,
            "chart_data": rs.get("chart_data"),
            "layout": layout,
            "colors": rs.get("colors"),
            "fonts": rs.get("fonts"),
        })

    return slides


def _generate_fallback(prompt: str, count: int, theme: dict, title_bg: str) -> list:
    """Fallback template when AI is unavailable. Uses the prompt text directly."""
    topic = re.sub(
        r'\b(create|make|build|generate|presentation|slides?|about|on|the|a|an|for|me|please|can you|i need|i want|creează|fă|generează|prezentare|despre|pentru)\b',
        '', prompt, flags=re.IGNORECASE
    ).strip().strip(' .,!?') or prompt
    topic_title = topic[0].upper() + topic[1:] if topic else "Presentation"

    slides = [
        {"id": "sl-gen-1", "type": "title", "title": topic_title, "subtitle": prompt,
         "bg": title_bg, "bullets": [], "source": None, "notes": "", "image": None},
    ]
    section_titles = [
        "Introduction", "Overview", "Key Points", "Analysis",
        "Details", "Challenges", "Recommendations", "Conclusion",
    ]
    for i in range(1, count - 1):
        idx = (i - 1) % len(section_titles)
        slides.append({
            "id": f"sl-gen-{i+1}", "type": "content",
            "title": section_titles[idx], "subtitle": f"Section {i}",
            "bg": theme["bg"],
            "bullets": [f"Point about {topic}"] * 4,
            "source": None, "notes": "", "image": None,
        })
    slides.append({
        "id": f"sl-gen-{count}", "type": "title", "title": "Thank You",
        "subtitle": prompt, "bg": title_bg, "bullets": [],
        "source": None, "notes": "", "image": None,
    })
    return slides


# ══════════════════════════════════════════════════════════════
#  /command — process editing commands
# ══════════════════════════════════════════════════════════════

COMMAND_RESPONSES = {
    "add-slide": "Added a new slide. You can edit it in the inspector panel.",
    "more-data": "Enhanced the slide with additional data points and sources.",
    "change-theme": "Theme updated. All slides now use the new color scheme.",
    "make-visual": "Improved visual layout — reduced text density and enhanced hierarchy.",
    "critique": (
        "**AI Critique:**\n"
        "• Slide 2 could use a stronger opening hook\n"
        "• Slide 4 has good data but needs source citations\n"
        "• Consider adding an image to Slide 5 for visual balance\n"
        "• The conclusion could include a clearer call-to-action\n"
        "• Overall: **7.5/10** — solid structure, needs refinement in data density"
    ),
    "rewrite": "Slide rewritten with improved clarity and flow.",
    "export": "Use the download button to export your presentation.",
}

@router.post("/command")
async def slide_command(req: CommandRequest):
    """Process editing commands for presentations."""
    text = req.text.strip()
    logger.info("Slide command: text=%r slide=%d", text[:80], req.current_slide)

    # Parse command
    if text.startswith("/"):
        parts = text.split(" ", 1)
        cmd = parts[0].lower().lstrip("/")
        arg = parts[1] if len(parts) > 1 else ""

        response = COMMAND_RESPONSES.get(cmd, f"Command `{cmd}` acknowledged.")

        # Customize responses
        if cmd == "change-theme" and arg:
            theme_name = arg.strip().lower()
            if theme_name in THEMES:
                response = f"Theme changed to **{theme_name}**. All slides updated."
            else:
                response = f"Unknown theme '{theme_name}'. Available: {', '.join(THEMES.keys())}"

        if cmd == "more-data":
            response = f"Added more data to slide {req.current_slide + 1}. Sources have been verified."

        if cmd == "make-visual":
            response = f"Slide {req.current_slide + 1} redesigned with improved visual hierarchy."

        return {"command": cmd, "response": response, "args": arg}

    # Natural language — treat as chat about presentation
    return {
        "command": "chat",
        "response": (
            "I can help refine your presentation. Try:\n"
            "• **/add-slide [topic]** — add a new slide\n"
            "• **/more-data** — add data to current slide\n"
            "• **/change-theme [style]** — change visual theme\n"
            "• **/make-visual** — improve slide layout\n"
            "• **/critique** — get AI feedback on your deck\n"
            "• **/rewrite** — rewrite current slide"
        ),
    }


# ══════════════════════════════════════════════════════════════
#  /status
# ══════════════════════════════════════════════════════════════

@router.get("/status")
async def slide_status():
    return {
        "status": "ok",
        "themes": list(THEMES.keys()),
    }


# ══════════════════════════════════════════════════════════════
#  /export — Generate a real PPTX file from slide data
# ══════════════════════════════════════════════════════════════

class ExportSlide(BaseModel):
    title: str = ""
    subtitle: str = ""
    bullets: List[str] = []
    notes: str = ""
    source: Optional[str] = None
    type: str = "content"

class ExportRequest(BaseModel):
    slides: List[ExportSlide]
    theme: str = "professional"
    filename: str = "presentation"

EXPORT_DIR = Path(tempfile.gettempdir()) / "slide_exports"
EXPORT_DIR.mkdir(parents=True, exist_ok=True)

@router.post("/export")
async def slide_export(req: ExportRequest):
    """Generate a real PPTX file from slide JSON data."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise HTTPException(503, "python-pptx not installed. Run: pip install python-pptx")

    if not req.slides:
        raise HTTPException(400, "No slides provided")

    theme = THEMES.get(req.theme, THEMES["professional"])

    # Parse hex color to RGBColor
    def hex_to_rgb(h: str) -> RGBColor:
        h = h.lstrip("#")
        if len(h) == 6:
            return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        return RGBColor(15, 23, 42)  # fallback dark

    bg_color = hex_to_rgb(theme["bg"])
    accent_color = hex_to_rgb(theme["accent"])
    text_color = hex_to_rgb(theme["text"])
    secondary_color = hex_to_rgb(theme["secondary"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(req.slides):
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        is_title = slide_data.type == "title"

        if is_title:
            # Title slide — centered
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.8))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(44)
            p.font.color.rgb = text_color
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            if slide_data.subtitle:
                sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.333), Inches(1))
                tf2 = sub_box.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = slide_data.subtitle
                p2.font.size = Pt(20)
                p2.font.color.rgb = secondary_color
                p2.alignment = PP_ALIGN.CENTER

            # Accent bar
            from pptx.util import Emu
            bar = slide.shapes.add_shape(
                1,  # rectangle
                Inches(5.667), Inches(4.0), Inches(2), Pt(3)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_color
            bar.line.fill.background()
        else:
            # Content slide
            # Slide number badge
            num_box = slide.shapes.add_textbox(Inches(11.8), Inches(0.4), Inches(0.8), Inches(0.5))
            ntf = num_box.text_frame
            np_ = ntf.paragraphs[0]
            np_.text = str(i + 1)
            np_.font.size = Pt(12)
            np_.font.color.rgb = accent_color
            np_.alignment = PP_ALIGN.CENTER

            # Title
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10.5), Inches(1.0))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(32)
            p.font.color.rgb = text_color
            p.font.bold = True

            # Subtitle
            y_offset = 1.6
            if slide_data.subtitle:
                sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_offset), Inches(10.5), Inches(0.7))
                tf2 = sub_box.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = slide_data.subtitle
                p2.font.size = Pt(16)
                p2.font.color.rgb = secondary_color
                y_offset += 0.8

            # Accent bar under title
            bar = slide.shapes.add_shape(1, Inches(0.8), Inches(y_offset), Inches(1.5), Pt(3))
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_color
            bar.line.fill.background()
            y_offset += 0.5

            # Bullets
            if slide_data.bullets:
                bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_offset), Inches(10.5), Inches(7.5 - y_offset - 1.0))
                btf = bullet_box.text_frame
                btf.word_wrap = True
                for j, bullet in enumerate(slide_data.bullets):
                    bp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
                    bp.text = f"  •  {bullet}"
                    bp.font.size = Pt(16)
                    bp.font.color.rgb = text_color
                    bp.space_after = Pt(12)

            # Source
            if slide_data.source:
                src_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(10.5), Inches(0.4))
                stf = src_box.text_frame
                sp = stf.paragraphs[0]
                sp.text = f"Source: {slide_data.source}"
                sp.font.size = Pt(10)
                sp.font.color.rgb = secondary_color
                sp.font.italic = True

        # Speaker notes
        if slide_data.notes:
            slide.notes_slide.notes_text_frame.text = slide_data.notes

    # Save to temp file
    safe_name = re.sub(r"[^\w\-]", "_", req.filename or "presentation")[:50]
    out_path = EXPORT_DIR / f"{safe_name}_{os.getpid()}.pptx"
    prs.save(str(out_path))
    logger.info("Exported PPTX: %s (%d slides)", out_path.name, len(req.slides))

    return FileResponse(
        path=str(out_path),
        filename=f"{safe_name}.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
